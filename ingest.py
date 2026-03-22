#!/usr/bin/env python3
"""
ingest.py - Scan the Public vault, extract text, convert to Markdown,
generate embeddings via Ollama, and store them in a local ChromaDB instance.

Usage:
    source /Users/chaomingou/Documents/VaultTools/.venv/bin/activate
    python /Users/chaomingou/Documents/VaultTools/ingest.py [subpath]
"""

import os
import hashlib
import json
import chromadb
from pathlib import Path
from typing import Optional, List

# ─── Configuration ───────────────────────────────────────────────────────────

VAULT_PATH = Path("/Users/chaomingou/Public")
CHROMA_PATH = Path("/Users/chaomingou/Documents/VaultTools/.chromadb")
COLLECTION_NAME = "public_vault"
OLLAMA_EMBED_MODEL = "nomic-embed-text"

# File types we support
SUPPORTED_EXTENSIONS = {".md", ".txt", ".pdf", ".xlsx", ".csv", ".pptx"}

# Directories to skip
SKIP_DIRS = {".venv", ".obsidian", ".chromadb", ".git", "node_modules", "99_Templates"}

# Chunk size (characters) — larger for better context
CHUNK_SIZE = 1500
CHUNK_OVERLAP = 300


# ─── Text Extraction ────────────────────────────────────────────────────────

def extract_text_from_md(filepath: Path) -> str:
    """Read markdown/text files."""
    try:
        return filepath.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print(f"  ⚠ Error reading {filepath.name}: {e}")
        return ""


def extract_text_from_pdf(filepath: Path) -> str:
    """Extract text from PDF using pypdf."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(filepath))
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        print(f"  ⚠ Error reading PDF {filepath.name}: {e}")
        return ""


def extract_text_from_xlsx(filepath: Path) -> str:
    """Extract text from Excel files."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(filepath), read_only=True, data_only=True)
        text_parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_parts.append(f"--- Sheet: {sheet} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        wb.close()
        return "\n".join(text_parts)
    except Exception as e:
        print(f"  ⚠ Error reading Excel {filepath.name}: {e}")
        return ""


def extract_text_from_csv(filepath: Path) -> str:
    """Read CSV files as text."""
    try:
        return filepath.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print(f"  ⚠ Error reading CSV {filepath.name}: {e}")
        return ""


def extract_text_from_pptx(filepath: Path) -> str:
    """Extract text from PowerPoint files using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        print(f"  ⚠ Error reading PPTX {filepath.name}: {e}")
        return ""


def extract_text(filepath: Path) -> str:
    """Route file to the correct extractor."""
    ext = filepath.suffix.lower()
    if ext in {".md", ".txt"}:
        return extract_text_from_md(filepath)
    elif ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(filepath)
    elif ext == ".csv":
        return extract_text_from_csv(filepath)
    elif ext == ".pptx":
        return extract_text_from_pptx(filepath)
    return ""


# ─── Markdown Conversion ────────────────────────────────────────────────────

def convert_to_markdown(text: str, filepath: Path) -> str:
    """Convert extracted text to a structured Markdown format for better chunking."""
    ext = filepath.suffix.lower()
    filename = filepath.stem
    category = detect_category(filepath)

    # Build a Markdown header with metadata
    header = f"# {filename}\n\n"
    header += f"> **File**: {filepath.name}  \n"
    header += f"> **Category**: {category}  \n\n"

    if ext == ".pdf":
        # For PDFs: add section separators where there are large gaps
        lines = text.split("\n")
        md_lines = []
        prev_empty = 0
        for line in lines:
            stripped = line.strip()
            if not stripped:
                prev_empty += 1
                if prev_empty == 3:
                    md_lines.append("\n---\n")
                continue
            prev_empty = 0
            md_lines.append(stripped)
        return header + "\n".join(md_lines)

    elif ext == ".xlsx":
        # For Excel: preserve sheet headers and format as tables
        lines = text.split("\n")
        md_lines = []
        for line in lines:
            if line.startswith("--- Sheet:"):
                sheet_name = line.replace("--- Sheet:", "").replace("---", "").strip()
                md_lines.append(f"\n## Sheet: {sheet_name}\n")
            elif "|" in line:
                md_lines.append(f"| {line} |")
            else:
                md_lines.append(line)
        return header + "\n".join(md_lines)

    elif ext == ".csv":
        # For CSV: convert to markdown table
        lines = text.strip().split("\n")
        if not lines:
            return header
        md_lines = [f"\n## Data\n"]
        for i, line in enumerate(lines):
            cells = line.split(",")
            md_lines.append("| " + " | ".join(c.strip() for c in cells) + " |")
            if i == 0:
                md_lines.append("| " + " | ".join("---" for _ in cells) + " |")
        return header + "\n".join(md_lines)

    elif ext == ".pptx":
        # For PPTX: separate slides
        lines = text.split("\n")
        md_lines = []
        slide_num = 1
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            if len(stripped) < 60 and not stripped[0].isdigit():
                md_lines.append(f"\n### Slide {slide_num}: {stripped}\n")
                slide_num += 1
            else:
                md_lines.append(stripped)
        return header + "\n".join(md_lines)

    else:
        # .md and .txt — already markdown-ish
        return header + text


# ─── Chunking ────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping chunks, trying to break at paragraph boundaries."""
    if not text.strip():
        return []

    chunks = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = start + chunk_size

        # Try to find a good break point (paragraph boundary)
        if end < text_len:
            # Look for double newline within last 200 chars
            break_zone = text[max(end - 200, start):end]
            last_para = break_zone.rfind("\n\n")
            if last_para != -1:
                end = max(end - 200, start) + last_para + 2
            else:
                # Try single newline
                last_nl = break_zone.rfind("\n")
                if last_nl != -1:
                    end = max(end - 200, start) + last_nl + 1

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        # Move forward, accounting for overlap
        start = end - overlap if end < text_len else text_len

    return chunks


# ─── Embedding via Ollama ────────────────────────────────────────────────────

def get_embedding(text: str) -> Optional[List[float]]:
    """Get embedding vector from Ollama."""
    prefixed_text = f"search_document: {text}"
    try:
        import urllib.request
        data = json.dumps({"model": OLLAMA_EMBED_MODEL, "input": prefixed_text[:2000]}).encode()
        req = urllib.request.Request(
            "http://localhost:11434/api/embed",
            data=data,
            headers={"Content-Type": "application/json"}
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read().decode())
            embeddings = result.get("embeddings")
            if embeddings and len(embeddings) > 0:
                return embeddings[0]
        return None
    except Exception as e:
        print(f"  ⚠ Embedding error: {e}")
        return None


# ─── File Discovery ──────────────────────────────────────────────────────────

def discover_files(vault_path: Path) -> list[Path]:
    """Find all supported files in the vault."""
    files = []
    for root, dirs, filenames in os.walk(vault_path):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS and not d.startswith(".")]
        for fname in filenames:
            fpath = Path(root) / fname
            if fpath.suffix.lower() in SUPPORTED_EXTENSIONS:
                files.append(fpath)
    return sorted(files)


def file_hash(filepath: Path) -> str:
    """Generate a hash for the file to detect changes."""
    stat = filepath.stat()
    return hashlib.md5(f"{filepath}:{stat.st_size}:{stat.st_mtime}".encode()).hexdigest()


# ─── Category Detection ─────────────────────────────────────────────────────

def detect_category(filepath: Path) -> str:
    """Detect the category based on the folder structure."""
    try:
        rel = filepath.relative_to(VAULT_PATH)
        parts = rel.parts
        if len(parts) > 1:
            return parts[0]
    except ValueError:
        pass
    return "Root"


# ─── Main Ingestion ──────────────────────────────────────────────────────────

def main():
    import sys
    subpath = sys.argv[1] if len(sys.argv) > 1 else ""
    target_vault = VAULT_PATH / subpath

    print("🧠 Public Vault Ingestion")
    print(f"   Vault: {target_vault}")
    print(f"   ChromaDB: {CHROMA_PATH}")
    print(f"   Chunk size: {CHUNK_SIZE} chars, overlap: {CHUNK_OVERLAP} chars")
    print(f"   Markdown conversion: enabled")
    print()

    # Initialize ChromaDB
    client = chromadb.PersistentClient(path=str(CHROMA_PATH))
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        metadata={"hnsw:space": "cosine"}
    )

    # Discover files
    files = discover_files(target_vault)
    print(f"📂 Found {len(files)} supported files\n")

    total_chunks = 0
    errors = 0

    for i, filepath in enumerate(files, 1):
        category = detect_category(filepath)
        fhash = file_hash(filepath)
        first_chunk_id = f"{fhash}_0"

        existing = collection.get(ids=[first_chunk_id], include=["metadatas"])
        if existing and existing["ids"]:
            print(f"[{i}/{len(files)}] ⏭ {category}/{filepath.name} (Already indexed, skipping)")
            continue

        print(f"[{i}/{len(files)}] 🔄 {category}/{filepath.name}")

        # Extract text
        text = extract_text(filepath)
        if not text.strip():
            print("  ⏭ No text extracted, skipping.")
            continue

        # Convert to Markdown
        md_text = convert_to_markdown(text, filepath)

        # Chunk the markdown text
        chunks = chunk_text(md_text)
        if not chunks:
            print("  ⏭ No chunks generated, skipping.")
            continue

        print(f"  📝 {len(chunks)} chunks (markdown)")

        for j, chunk in enumerate(chunks):
            chunk_id = f"{fhash}_{j}"

            # Get embedding
            embedding = get_embedding(chunk)
            if embedding is None:
                errors += 1
                continue

            # Store in ChromaDB with file_hash metadata
            collection.add(
                ids=[chunk_id],
                embeddings=[embedding],
                documents=[chunk],
                metadatas=[{
                    "source": str(filepath),
                    "filename": filepath.name,
                    "category": category,
                    "chunk_index": j,
                    "total_chunks": len(chunks),
                    "file_hash": fhash,
                }]
            )
            total_chunks += 1

        print(f"  ✅ Stored {len(chunks)} chunks")

    print(f"\n{'='*50}")
    print(f"✅ Ingestion complete!")
    print(f"   Total files processed: {len(files)}")
    print(f"   Total chunks stored: {total_chunks}")
    print(f"   Errors: {errors}")
    print(f"   Database location: {CHROMA_PATH}")


if __name__ == "__main__":
    main()
